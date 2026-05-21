from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ── Colour palette ──────────────────────────────────────────────────────────
NAVY       = RGBColor(0x0D, 0x1B, 0x2A)   # slide background / header fill
ICE_BLUE   = RGBColor(0x00, 0x8B, 0xB8)   # accent / underlines
LIGHT_BG   = RGBColor(0xF0, 0xF4, 0xF8)   # body background
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT  = RGBColor(0x1A, 0x2B, 0x3C)
ACCENT     = RGBColor(0x00, 0xC6, 0xE0)   # highlight bullets

IMG_DIR = os.path.join(os.path.dirname(__file__), "slides_images")

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]   # completely blank

# ─────────────────────────────────────────────────────────────────────────────
# Helper utilities
# ─────────────────────────────────────────────────────────────────────────────

def add_rect(slide, l, t, w, h, fill_rgb=None, line_rgb=None, line_width_pt=0):
    from pptx.util import Pt as _Pt
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.fill.background()
    if fill_rgb:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_rgb
    else:
        shape.fill.background()
    if line_rgb:
        shape.line.color.rgb = line_rgb
        shape.line.width = Pt(line_width_pt) if line_width_pt else Pt(0)
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, text, l, t, w, h,
                font_size=18, bold=False, italic=False,
                color=DARK_TEXT, align=PP_ALIGN.LEFT,
                wrap=True, font_name="Calibri"):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_name
    return txb


def add_bullet_textbox(slide, bullets, l, t, w, h,
                       font_size=16, color=DARK_TEXT, font_name="Calibri",
                       bullet_char="•"):
    from pptx.util import Pt as _Pt
    from pptx.oxml.ns import qn
    from lxml import etree
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for item in bullets:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = f"{bullet_char}  {item}"
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.name = font_name
    return txb


def add_image(slide, img_path, l, t, w, h=None):
    if h is None:
        pic = slide.shapes.add_picture(img_path, Inches(l), Inches(t), width=Inches(w))
    else:
        pic = slide.shapes.add_picture(img_path, Inches(l), Inches(t), width=Inches(w), height=Inches(h))
    return pic


def add_slide_background(slide, color=LIGHT_BG):
    """Fill the whole slide with a solid colour."""
    add_rect(slide, 0, 0, 13.33, 7.5, fill_rgb=color)


def add_header_bar(slide, title_text, subtitle_text=""):
    """Dark banner at top + title text."""
    add_rect(slide, 0, 0, 13.33, 1.35, fill_rgb=NAVY)
    add_rect(slide, 0, 1.35, 13.33, 0.045, fill_rgb=ICE_BLUE)   # accent line
    add_textbox(slide, title_text, 0.35, 0.08, 12.5, 0.8,
                font_size=28, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle_text:
        add_textbox(slide, subtitle_text, 0.35, 0.82, 12.5, 0.45,
                    font_size=14, bold=False, color=ACCENT, align=PP_ALIGN.LEFT)


def add_section_label(slide, text, l=0.35, t=1.55, w=12.6, font_size=13, color=ICE_BLUE):
    add_textbox(slide, text.upper(), l, t, w, 0.3,
                font_size=font_size, bold=True, color=color)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 1 — Title
# ─────────────────────────────────────────────────────────────────────────────
def slide_title():
    slide = prs.slides.add_slide(BLANK)
    # Full-bleed dark background
    add_rect(slide, 0, 0, 13.33, 7.5, fill_rgb=NAVY)
    # Accent stripe
    add_rect(slide, 0, 4.8, 13.33, 0.07, fill_rgb=ICE_BLUE)
    # Group info box
    add_rect(slide, 0, 4.87, 13.33, 2.63, fill_rgb=RGBColor(0x08, 0x12, 0x1C))

    # Title
    add_textbox(slide, "Antarctic Microbial Sensor", 1.0, 1.3, 11.0, 1.1,
                font_size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, "Autonomous Marine Environmental Monitoring System", 1.0, 2.5, 11.0, 0.65,
                font_size=22, bold=False, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, "EEE4113F  ·  Group 14  ·  University of Cape Town", 1.0, 3.3, 11.0, 0.5,
                font_size=16, bold=False, color=RGBColor(0xA0, 0xC4, 0xD8), align=PP_ALIGN.CENTER)

    # Team
    add_textbox(slide, "Devon Clark   ·   Joshua Naidoo   ·   Kiyuran Naidoo   ·   Saeed Solomon",
                1.0, 5.1, 11.0, 0.5,
                font_size=15, bold=False, color=RGBColor(0xB0, 0xC8, 0xDC), align=PP_ALIGN.CENTER)
    add_textbox(slide, "Supervised by Dr Emma Rocke  ·  MARiS, UCT",
                1.0, 5.7, 11.0, 0.4,
                font_size=13, bold=False, color=RGBColor(0x70, 0x90, 0xA8), align=PP_ALIGN.CENTER)
    add_textbox(slide, "May 2026",
                1.0, 6.2, 11.0, 0.4,
                font_size=13, bold=False, color=RGBColor(0x70, 0x90, 0xA8), align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 2 — Background & Motivation
# ─────────────────────────────────────────────────────────────────────────────
def slide_background():
    slide = prs.slides.add_slide(BLANK)
    add_slide_background(slide)
    add_header_bar(slide, "Background & Motivation",
                   "Why does Antarctic marine monitoring matter?")

    bullets_l = [
        "Marine microbes are the 'invisible engines' of the ocean ecosystem",
        "They are the first biological indicators of climate change and pollution",
        "Dr Emma Rocke (MARiS, UCT) studies microbial dynamics across Antarctic & Southern Oceans",
        "Early detection of harmful algal blooms requires continuous, high-resolution data",
    ]
    bullets_r = [
        "Current practice relies on expensive, infrequent research cruises",
        "Austral winter (-30 °C, expanding sea ice) makes the region inaccessible for months",
        "This creates a critical 'winter data gap' — the time most valuable for early warnings",
        "Human presence in the Antarctic is hazardous, costly, and logistically prohibitive",
    ]

    # Left panel
    add_rect(slide, 0.3, 1.55, 6.1, 5.6, fill_rgb=WHITE,
             line_rgb=ICE_BLUE, line_width_pt=1.2)
    add_textbox(slide, "THE OCEAN'S EARLY-WARNING SYSTEM", 0.5, 1.65, 5.7, 0.4,
                font_size=11, bold=True, color=ICE_BLUE)
    add_bullet_textbox(slide, bullets_l, 0.5, 2.05, 5.7, 5.0, font_size=15)

    # Right panel
    add_rect(slide, 6.9, 1.55, 6.1, 5.6, fill_rgb=WHITE,
             line_rgb=RGBColor(0xE0, 0x44, 0x44), line_width_pt=1.2)
    add_textbox(slide, "THE 'WINTER DATA GAP'", 7.1, 1.65, 5.7, 0.4,
                font_size=11, bold=True, color=RGBColor(0xC0, 0x30, 0x30))
    add_bullet_textbox(slide, bullets_r, 7.1, 2.05, 5.7, 5.0, font_size=15)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 3 — Project Objectives
# ─────────────────────────────────────────────────────────────────────────────
def slide_objectives():
    slide = prs.slides.add_slide(BLANK)
    add_slide_background(slide)
    add_header_bar(slide, "Project Objectives",
                   "Design a low-cost, autonomous bipartite marine monitoring system")

    objectives = [
        ("Continuous Sensing",
         "Acquire chlorophyll-a fluorescence, temperature, and depth as microbial community proxies"),
        ("Reliable Storage",
         "Store sensor data in non-volatile memory under extreme sub-zero conditions during a dive"),
        ("Wireless Offload",
         "Transfer stored data wirelessly to a researcher's laptop upon probe retrieval — no physical port"),
        ("Offline Intelligence",
         "Process data locally using ML anomaly detection to identify ecological events — fully offline"),
        ("Annual Autonomy",
         "Sustain operation for up to 12 months on a single primary battery without human maintenance"),
        ("Low Cost & Safe",
         "Total component cost ≤ R2 000; hand-deployable by one non-specialist from a vessel of opportunity"),
    ]

    cols = [(0.3, 4.15), (6.85, 4.15)]
    row_h = 1.55
    for i, (title, desc) in enumerate(objectives):
        col = i % 2
        row = i // 2
        lx = cols[col][0]
        ty = 1.55 + row * row_h
        add_rect(slide, lx, ty + 0.05, cols[col][1], row_h - 0.15,
                 fill_rgb=WHITE, line_rgb=ICE_BLUE, line_width_pt=1)
        add_rect(slide, lx, ty + 0.05, cols[col][1], 0.38, fill_rgb=NAVY)
        add_textbox(slide, f"{'①②③④⑤⑥'[i]}  {title}",
                    lx + 0.15, ty + 0.1, cols[col][1] - 0.3, 0.35,
                    font_size=13, bold=True, color=WHITE)
        add_textbox(slide, desc,
                    lx + 0.15, ty + 0.5, cols[col][1] - 0.3, row_h - 0.7,
                    font_size=13, color=DARK_TEXT)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 4 — System Architecture Overview
# ─────────────────────────────────────────────────────────────────────────────
def slide_system_overview():
    slide = prs.slides.add_slide(BLANK)
    add_slide_background(slide)
    add_header_bar(slide, "System Architecture Overview",
                   "A bipartite store-and-forward system deployed from vessels of opportunity")

    # Central description
    add_textbox(slide,
        "The system comprises two hardware nodes and one software platform. "
        "The submersible probe collects data autonomously underwater, then offloads "
        "wirelessly to a USB dongle connected to the researcher's laptop.",
        0.35, 1.55, 12.6, 0.7, font_size=15, color=DARK_TEXT)

    # Three boxes: Probe | Wireless | Laptop
    boxes = [
        (0.3,  2.4, 3.8, 4.7, NAVY,  WHITE,  "SUBMERSIBLE PROBE",
         ["ESP32-C6 (FireBeetle 2)", "FRAM non-volatile storage",
          "AS7341 spectral sensor", "DS18B20 temperature", "Pressure/depth sensor",
          "Li-SOCl₂ primary battery", "Thermal insulation shell"]),
        (4.75, 2.4, 3.8, 4.7, ICE_BLUE, WHITE, "ESP-NOW WIRELESS LINK",
         ["Connectionless data-link protocol", "100 m range; ultra-low power",
          "CRC-16 packet integrity check", "Reed switch wakeup from deep sleep",
          "ACK/NACK retry mechanism", "< 60 s full data offload"]),
        (9.2,  2.4, 3.8, 4.7, NAVY,  WHITE, "MISSION CONTROL (LAPTOP)",
         ["Receiver dongle (ESP32-C6 via USB)", "Python GUI (CustomTkinter)",
          "Raw data table + interactive graphs", "Isolation Forest + XAI anomaly detection",
          "Save / load session files", "100 % offline — no internet required"]),
    ]

    for lx, ty, bw, bh, hdr_col, txt_col, title, items in boxes:
        add_rect(slide, lx, ty, bw, bh, fill_rgb=WHITE,
                 line_rgb=ICE_BLUE, line_width_pt=1.2)
        add_rect(slide, lx, ty, bw, 0.45, fill_rgb=hdr_col)
        add_textbox(slide, title, lx + 0.1, ty + 0.07, bw - 0.2, 0.38,
                    font_size=11, bold=True, color=txt_col, align=PP_ALIGN.CENTER)
        add_bullet_textbox(slide, items, lx + 0.12, ty + 0.52, bw - 0.24, bh - 0.6,
                           font_size=12.5, color=DARK_TEXT)

    # Arrows between boxes
    add_textbox(slide, "⇌  ESP-NOW", 4.2, 4.5, 0.5, 0.4,
                font_size=10, color=DARK_TEXT, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 5 — Data Transfer & Visualisation
# ─────────────────────────────────────────────────────────────────────────────
def slide_data_transfer():
    slide = prs.slides.add_slide(BLANK)
    add_slide_background(slide)
    add_header_bar(slide, "Subsystem 1 — Data Transfer & Visualisation",
                   "Saeed Solomon  ·  ESP-NOW wireless protocol + Mission Control GUI")

    # Left: key design choices as text
    add_rect(slide, 0.3, 1.55, 5.0, 5.6, fill_rgb=WHITE,
             line_rgb=ICE_BLUE, line_width_pt=1)
    add_textbox(slide, "PROTOCOL SELECTION", 0.5, 1.65, 4.6, 0.35,
                font_size=11, bold=True, color=ICE_BLUE)

    proto_items = [
        "NFC — ruled out: strict proximity & orientation requirements",
        "BLE — implemented but abandoned: ESP32-C6 stack immaturity",
        "ESP-NOW — selected: proven reliability, 100 m range, ultra-low power",
    ]
    add_bullet_textbox(slide, proto_items, 0.5, 2.05, 4.6, 1.3,
                       font_size=13, color=DARK_TEXT)

    add_textbox(slide, "PROBE COMMAND FLOW", 0.5, 3.45, 4.6, 0.35,
                font_size=11, bold=True, color=ICE_BLUE)
    cmd_items = [
        "[CMD]:CONNECT — verify link, update battery level on GUI",
        "[CMD]:PREPARE — reset FRAM, store session timestamp",
        "[CMD]:RETRIEVE — stream all records with CRC-16 check",
        "[CMD]:DISCONNECT — power down radio to conserve battery",
    ]
    add_bullet_textbox(slide, cmd_items, 0.5, 3.85, 4.6, 2.0,
                       font_size=13, color=DARK_TEXT)

    add_textbox(slide, "RESULTS", 0.5, 5.9, 4.6, 0.35,
                font_size=11, bold=True, color=ICE_BLUE)
    results = ["1.2 % avg data error (spec: < 5 %)",
               "8 s wake-to-advertising (spec: < 10 s)",
               "All 7 ATPs passed"]
    add_bullet_textbox(slide, results, 0.5, 6.28, 4.6, 1.0, font_size=13, color=DARK_TEXT)

    # Right top: block diagram
    img_block = os.path.join(IMG_DIR, "page25_img1.jpeg")
    add_image(slide, img_block, 5.55, 1.55, 7.45, 2.7)

    # Right bottom: state-flow diagram
    img_flow = os.path.join(IMG_DIR, "page28_img1.jpeg")
    add_image(slide, img_flow, 5.55, 4.35, 7.45, 2.8)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 6 — Mission Control GUI
# ─────────────────────────────────────────────────────────────────────────────
def slide_gui():
    slide = prs.slides.add_slide(BLANK)
    add_slide_background(slide)
    add_header_bar(slide, "Mission Control GUI",
                   "Fully offline Python desktop application — CustomTkinter + Matplotlib")

    # Description strip
    desc_items = [
        "Top toolbar shows connection status and battery percentage",
        "Sidebar: Connect / Prepare Dive / Retrieve Data / Dev Mode / Save & Load",
        "Raw Data tab — scrollable table of all 14 sensor channels per record",
        "Graph tab — interactive time-series with All Samples or per-Group view",
        "Status log at the bottom prints timestamped system events",
        "Anomaly markers (red dots) overlaid on graphs by the ML backend",
    ]
    add_rect(slide, 0.3, 1.55, 5.0, 5.6, fill_rgb=WHITE,
             line_rgb=ICE_BLUE, line_width_pt=1)
    add_textbox(slide, "KEY FEATURES", 0.5, 1.65, 4.6, 0.35,
                font_size=11, bold=True, color=ICE_BLUE)
    add_bullet_textbox(slide, desc_items, 0.5, 2.05, 4.6, 5.0,
                       font_size=13.5, color=DARK_TEXT)

    # Two GUI screenshots
    img_raw  = os.path.join(IMG_DIR, "page30_img1.png")
    img_graph = os.path.join(IMG_DIR, "page30_img2.png")
    add_image(slide, img_raw,   5.55, 1.55, 7.5, 2.65)
    add_image(slide, img_graph, 5.55, 4.3,  7.5, 2.8)

    add_textbox(slide, "Raw Data tab", 5.55, 4.22, 3.6, 0.25,
                font_size=10, color=ICE_BLUE, bold=True)
    add_textbox(slide, "Graph tab", 9.6, 4.22, 3.4, 0.25,
                font_size=10, color=ICE_BLUE, bold=True)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 7 — Edge Architecture & Explainable AI
# ─────────────────────────────────────────────────────────────────────────────
def slide_edge_ai():
    slide = prs.slides.add_slide(BLANK)
    add_slide_background(slide)
    add_header_bar(slide, "Subsystem 2 — Edge Architecture & Explainable AI",
                   "Kiyuran Naidoo  ·  On-probe firmware lifecycle + offline ML anomaly detection")

    # Left: firmware states
    add_rect(slide, 0.3, 1.55, 5.8, 2.8, fill_rgb=WHITE,
             line_rgb=ICE_BLUE, line_width_pt=1)
    add_textbox(slide, "4-STATE FIRMWARE LIFECYCLE", 0.5, 1.65, 5.4, 0.35,
                font_size=11, bold=True, color=ICE_BLUE)
    states = [
        "STATE_IDLE — deep sleep on LP-IO, woken by magnetic reed switch",
        "STATE_ARMED — sinking delay (15 min) before data acquisition begins",
        "STATE_LOGGING — periodic measurement cycles; records written to FRAM",
        "STATE_OFFLOAD — activates ESP-NOW; awaits [CMD]:RETRIEVE from GUI",
    ]
    add_bullet_textbox(slide, states, 0.5, 2.06, 5.5, 2.2, font_size=13.5)

    # Left: FRAM storage
    add_rect(slide, 0.3, 4.45, 5.8, 1.4, fill_rgb=WHITE,
             line_rgb=ICE_BLUE, line_width_pt=1)
    add_textbox(slide, "MEMORY — FERROELECTRIC RAM (FRAM)", 0.5, 4.55, 5.4, 0.35,
                font_size=11, bold=True, color=ICE_BLUE)
    fram_items = [
        "10 mA write current vs 100 mA for SD card",
        "Operating range: -55 °C to +85 °C",
        "Data preserved across deep sleep resets (RTC SRAM)",
        "58-byte fixed MTU per record, 14 sensor values",
    ]
    add_bullet_textbox(slide, fram_items, 0.5, 4.92, 5.5, 1.0, font_size=13)

    # Left: XAI engine
    add_rect(slide, 0.3, 5.95, 5.8, 1.25, fill_rgb=WHITE,
             line_rgb=ICE_BLUE, line_width_pt=1)
    add_textbox(slide, "EXPLAINABLE AI (XAI) ENGINE", 0.5, 6.05, 5.4, 0.35,
                font_size=11, bold=True, color=ICE_BLUE)
    xai_items = [
        "Isolation Forest: anomaly isolation via random decision trees — O(n log n)",
        "Post-hoc deterministic ruleset: classifies Phytoplankton / Cyanobacteria / CDOM",
        "92 % accuracy, 0 false positives, ~0.008 s inference (spec: < 2 s)",
    ]
    add_bullet_textbox(slide, xai_items, 0.5, 6.42, 5.5, 0.75, font_size=12.5)

    # Right: anomaly graph
    img_anomaly = os.path.join(IMG_DIR, "page71_img1.png")
    add_image(slide, img_anomaly, 6.35, 1.55, 6.65, 4.1)
    add_textbox(slide,
        "Mission Control graph: 113 records with anomalies flagged in red by the MarineAnomalyDetector",
        6.35, 5.7, 6.65, 0.65, font_size=12, color=DARK_TEXT, italic=True)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 8 — Power & Thermal
# ─────────────────────────────────────────────────────────────────────────────
def slide_power():
    slide = prs.slides.add_slide(BLANK)
    add_slide_background(slide)
    add_header_bar(slide, "Subsystem 3 — Power & Thermal Management",
                   "Joshua Naidoo  ·  365-day autonomy in Antarctic cryogenic conditions")

    # Left: design choices
    add_rect(slide, 0.3, 1.55, 5.6, 5.6, fill_rgb=WHITE,
             line_rgb=ICE_BLUE, line_width_pt=1)
    add_textbox(slide, "ENERGY ARCHITECTURE", 0.5, 1.65, 5.2, 0.35,
                font_size=11, bold=True, color=ICE_BLUE)
    energy_items = [
        "ER14505 Li-SOCl₂ battery — operational down to -55 °C",
        "Nominal 2 400 mAh; 50 % capacity derated for -50 °C",
        "MCP1700 LDO: 1.6 µA quiescent vs 5 mA for LM7805",
        "FireBeetle 2 ESP32-C6: 10 µA deep sleep vs 10 mA (DevKit)",
        "Split-rail: 3.3 V always-on rail + 5 V switched sensor rail",
        "P-MOSFET (AO3401A) power gate — GPIO-controlled sensor isolation",
        "~22 measurement cycles/day over 365-day mission verified",
    ]
    add_bullet_textbox(slide, energy_items, 0.5, 2.06, 5.2, 2.4, font_size=13.5)

    add_textbox(slide, "THERMAL 'HUDDLE' STRATEGY", 0.5, 4.55, 5.2, 0.35,
                font_size=11, bold=True, color=ICE_BLUE)
    thermal_items = [
        "Layer 1 — Aluminium foil tape: reflects infrared radiation",
        "Layer 2 — EWELDH010 foam: stagnant air pockets, high thermal resistance",
        "RF window: 20 mm foil gap allows 2.4 GHz ESP-NOW signal to escape",
        "MCP1700 regulators thermally bonded to battery stack (self-heating)",
        "Result: internal core held at -12 °C after 6 h cold-soak at -30 °C",
        "16.8 µA measured standby current — within 20 µA specification",
    ]
    add_bullet_textbox(slide, thermal_items, 0.5, 4.95, 5.2, 2.25, font_size=13.5)

    # Right: circuit schematic
    img_circuit = os.path.join(IMG_DIR, "page49_img1.png")
    add_image(slide, img_circuit, 6.15, 1.55, 6.85, 5.6)
    add_textbox(slide,
        "Split-rail power schematic: ER14505 battery → MCP1700 LDOs → ESP32-C6 + FRAM + Sensors",
        6.15, 7.15, 6.85, 0.3, font_size=11, color=DARK_TEXT, italic=True)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 9 — Sensors & Housing
# ─────────────────────────────────────────────────────────────────────────────
def slide_sensors():
    slide = prs.slides.add_slide(BLANK)
    add_slide_background(slide)
    add_header_bar(slide, "Subsystem 4 — Sensors & Housing",
                   "Devon Clark  ·  Chlorophyll fluorometer, temperature, depth, waterproof PVC enclosure")

    # Left text
    add_rect(slide, 0.3, 1.55, 5.6, 5.6, fill_rgb=WHITE,
             line_rgb=ICE_BLUE, line_width_pt=1)
    add_textbox(slide, "SENSOR SUITE", 0.5, 1.65, 5.2, 0.35,
                font_size=11, bold=True, color=ICE_BLUE)
    sensor_items = [
        "AS7341 — 11-channel spectral sensor; 445 nm excitation, 680 nm fluorescence channels",
        "CZ104BC — 460 nm blue LED excitation (5 500–6 000 mcd); matches chlorophyll-a peak",
        "DS18B20 — waterproof 1-Wire digital; factory-calibrated; ±0.4 °C achieved",
        "Analog pressure sensor — 0–120 m range; R² > 0.99 depth linearity",
        "11-step PWM ramp algorithm to optimise fluorometer sensitivity",
    ]
    add_bullet_textbox(slide, sensor_items, 0.5, 2.06, 5.2, 2.3, font_size=13.5)

    add_textbox(slide, "WATERPROOF HOUSING", 0.5, 4.45, 5.2, 0.35,
                font_size=11, bold=True, color=ICE_BLUE)
    housing_items = [
        "50 mm OD PVC pipe main body (46.8 mm ID, 250 mm length)",
        "FDM 3D-printed PLA end caps with 4 perimeters (SLA recommended for production)",
        "Friction-fit double O-ring seal — field serviceable, no tools required",
        "Negatively buoyant — self-sinking with ballast weights",
        "Waterproof to 1 m for 30 min (prototype validation)",
        "Total system cost: < R2 000  ·  Mass: < 1 kg",
    ]
    add_bullet_textbox(slide, housing_items, 0.5, 4.85, 5.2, 2.3, font_size=13.5)

    # Right: two CAD renders side by side
    img_h1 = os.path.join(IMG_DIR, "page58_img1.png")
    img_h2 = os.path.join(IMG_DIR, "page58_img2.png")
    add_image(slide, img_h1, 6.1, 1.55, 3.55, 5.6)
    add_image(slide, img_h2, 9.75, 1.55, 3.35, 5.6)

    add_textbox(slide, "Rear end cap (sealed side)",
                6.1, 7.15, 3.55, 0.3, font_size=11, color=DARK_TEXT, italic=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, "Front end cap (sensor window side)",
                9.75, 7.15, 3.35, 0.3, font_size=11, color=DARK_TEXT, italic=True, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 10 — Results & Validation
# ─────────────────────────────────────────────────────────────────────────────
def slide_results():
    slide = prs.slides.add_slide(BLANK)
    add_slide_background(slide)
    add_header_bar(slide, "System Results & Validation",
                   "All subsystem acceptance test procedures passed")

    subsystems = [
        ("Data Transfer\n& Visualisation",
         ["1.2 % avg data error  (spec: < 5 %)",
          "8 s wake-to-advertising  (spec: < 10 s)",
          "100 m operational range  (spec: > 10 m)",
          "All 7 ATPs passed",
          "Fully offline GUI operation confirmed"]),
        ("Edge Architecture\n& Explainable AI",
         ["FRAM data intact across all deep sleep cycles",
          "4-state lifecycle validated via HITL serial telemetry",
          "92 % XAI classification accuracy, 0 false positives",
          "0.008 s inference per record  (spec: < 2 s)",
          "All 4 ATPs passed"]),
        ("Power & Thermal\nManagement",
         ["16.8 µA measured standby  (spec: ≤ 20 µA)",
          "-12 °C core after 6 h cold-soak at -30 °C",
          "3.3 V & 5.0 V rails stable across full discharge",
          "P-MOSFET gate: sensor current drops to 0 µA on GPIO",
          "All 5 ATPs passed"]),
        ("Sensors\n& Housing",
         ["Fluorescence ratio monotonically increases down to 1/32 dilution",
          "Temperature accuracy ±0.4 °C  (spec: ±0.5 °C)",
          "Depth linearity R² > 0.99  (spec: R² > 0.99)",
          "Waterproof at 1 m for 30 min — no ingress",
          "All 8 ATPs passed  ·  Total cost < R2 000"]),
    ]

    col_w = 3.1
    for i, (title, items) in enumerate(subsystems):
        lx = 0.27 + i * (col_w + 0.08)
        add_rect(slide, lx, 1.55, col_w, 5.7, fill_rgb=WHITE,
                 line_rgb=ICE_BLUE, line_width_pt=1)
        add_rect(slide, lx, 1.55, col_w, 0.6, fill_rgb=NAVY)
        add_textbox(slide, title, lx + 0.1, 1.58, col_w - 0.2, 0.55,
                    font_size=11.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_bullet_textbox(slide, items, lx + 0.12, 2.22, col_w - 0.24, 5.0,
                           font_size=12.5, color=DARK_TEXT)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 11 — Conclusions & Future Work
# ─────────────────────────────────────────────────────────────────────────────
def slide_conclusions():
    slide = prs.slides.add_slide(BLANK)
    add_slide_background(slide)
    add_header_bar(slide, "Conclusions & Future Work",
                   "A functional, integrated Antarctic monitoring system demonstrated across all four subsystems")

    # Conclusions
    add_rect(slide, 0.3, 1.55, 6.1, 5.65, fill_rgb=WHITE,
             line_rgb=ICE_BLUE, line_width_pt=1)
    add_textbox(slide, "CONCLUSIONS", 0.5, 1.65, 5.7, 0.35,
                font_size=11, bold=True, color=ICE_BLUE)
    conc = [
        "A fully integrated bipartite probe system was designed and demonstrated",
        "Sensor data is acquired, stored in FRAM, and wirelessly offloaded to a laptop GUI",
        "ML-based anomaly detection operates entirely offline on standard hardware",
        "System is hand-deployable by one non-specialist from a vessel of opportunity",
        "Annual autonomy validated through energy modelling and cold-soak testing",
        "Total cost within R2 000 budget; fits within 50 mm OD PVC housing",
        "All 24 acceptance test procedures across four subsystems passed",
    ]
    add_bullet_textbox(slide, conc, 0.5, 2.06, 5.7, 5.1, font_size=14, color=DARK_TEXT)

    # Future work
    add_rect(slide, 6.9, 1.55, 6.1, 5.65, fill_rgb=WHITE,
             line_rgb=RGBColor(0x00, 0x95, 0x50), line_width_pt=1)
    add_textbox(slide, "FUTURE WORK", 7.1, 1.65, 5.7, 0.35,
                font_size=11, bold=True, color=RGBColor(0x00, 0x85, 0x45))
    future = [
        "Migrate to BLE once ESP32-C6 NimBLE stack matures (remove dongle requirement)",
        "SLA resin end caps for long-term immersion (replace FDM PLA)",
        "Absolute chlorophyll calibration using field samples",
        "Extend depth rating beyond 2 m prototype limit",
        "Add dissolved oxygen & nutrient sensors within revised budget",
        "LoRaWAN / NVIS relay from docking station to mainland research facility",
        "Satellite uplink for real-time data access during deployment",
    ]
    add_bullet_textbox(slide, future, 7.1, 2.06, 5.7, 5.1, font_size=14, color=DARK_TEXT)


# ─────────────────────────────────────────────────────────────────────────────
# BUILD ALL SLIDES
# ─────────────────────────────────────────────────────────────────────────────
slide_title()
slide_background()
slide_objectives()
slide_system_overview()
slide_data_transfer()
slide_gui()
slide_edge_ai()
slide_power()
slide_sensors()
slide_results()
slide_conclusions()

out_path = os.path.join(os.path.dirname(__file__), "Antarctic_Probe_Slides.pptx")
prs.save(out_path)
print(f"Saved: {out_path}")
